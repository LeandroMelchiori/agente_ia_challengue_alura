"""
Genera la base de conocimiento ficticia de TechRetail Solutions S.R.L.
en los 8 formatos que pide el desafío (PDF, Word, Excel, PowerPoint,
Markdown, CSV, JSON, HTML), repartida en categorías de negocio.

Todo el contenido es ficticio pero coherente entre sí (misma empresa,
mismos planes, mismas integraciones), derivado del caso TechRetail.

Uso:
    python -m scripts.generar_documentos
"""
from __future__ import annotations

import csv
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from src.config import settings  # noqa: E402

DOCS = settings.docs_dir


def _dir(categoria: str) -> Path:
    d = DOCS / categoria
    d.mkdir(parents=True, exist_ok=True)
    return d


# ==========================================================================
#  MARKDOWN
# ==========================================================================
def gen_markdown():
    (_dir("estrategico") / "alcance_proyecto.md").write_text(
        """# TechRetail Solutions S.R.L. — Alcance del Proyecto

## Objetivo general
Desarrollar una plataforma SaaS de e-commerce para PyMEs y emprendedores
digitales en Argentina, que unifique catálogo, cobros, facturación fiscal y
logística en un solo sistema.

## Dentro del alcance
- Onboarding de comercios y carga de catálogo.
- Cobros con MercadoPago y Stripe.
- Facturación electrónica automática contra ARCA (WSFEv1).
- Motor de conciliación financiera automatizado.
- Conectores con ERPs locales (Tango Gestión y Bejerman).

## Fuera del alcance
- Punto de venta físico (POS) propio.
- Logística de última milla propia (se integra con Andreani, OCA y Correo Argentino).
- Operación fuera de Argentina en la primera etapa.

## Indicadores de control
- Tasa de activación de comercios en el trial.
- Churn mensual.
- MRR (Monthly Recurring Revenue).
- Discrepancia de conciliación (objetivo: < 0,1%).
""",
        encoding="utf-8",
    )

    (_dir("rrhh") / "manual_onboarding_colaboradores.md").write_text(
        """# Manual de Onboarding de Colaboradores — TechRetail

¡Bienvenido/a a TechRetail! Este manual resume tus primeros pasos.

## Primer día
1. Retirá tu notebook y credenciales de acceso con el equipo de IT.
2. Configurá el doble factor de autenticación (2FA) en todas las cuentas.
3. Sumate a los canales de Slack de tu squad.

## Estructura de squads
Trabajamos en squads multidisciplinarios (Producto, Pagos, Logística y
Plataforma). Cada squad tiene un Product Owner, un Tech Lead y desarrolladores.

## Modalidad de trabajo
- Esquema híbrido: 3 días remotos y 2 presenciales por semana.
- Horario flexible con núcleo obligatorio de 10 a 16 h.

## Herramientas
- Gestión de tareas: Jira.
- Documentación: Confluence.
- Comunicación: Slack.
- Repositorios: GitHub.
""",
        encoding="utf-8",
    )

    (_dir("financiero") / "politica_gastos_y_reembolsos.md").write_text(
        """# Política de Gastos y Reembolsos — TechRetail

## Alcance
Aplica a todos los colaboradores que incurran en gastos en nombre de la empresa.

## Gastos reembolsables
- Viáticos y transporte para reuniones con clientes.
- Herramientas de software con aprobación previa del líder de squad.
- Capacitaciones y cursos vinculados al rol.

## Gastos NO reembolsables
- Multas de tránsito.
- Consumos personales no relacionados con la actividad.

## Procedimiento
1. Cargá el comprobante (factura A o B a nombre de TechRetail Solutions S.R.L.)
   en el sistema de gastos dentro de los 30 días.
2. El líder de squad aprueba o rechaza en un plazo de 5 días hábiles.
3. Finanzas reintegra en el siguiente ciclo de pagos (los días 15 y 30).

## Topes
- Gastos superiores a ARS 200.000 requieren aprobación del área de Finanzas.
""",
        encoding="utf-8",
    )

    (_dir("legal") / "politica_privacidad.md").write_text(
        """# Política de Privacidad — TechRetail

TechRetail Solutions S.R.L. trata los datos personales conforme a la Ley
25.326 de Protección de Datos Personales de la República Argentina.

## Datos que recolectamos
- Datos de registro del comercio (razón social, CUIT, datos de contacto).
- Credenciales de integración fiscal y de pasarelas de pago.
- Datos de compradores finales necesarios para procesar transacciones.

## Finalidad
Los datos se usan exclusivamente para prestar el servicio: cobros,
facturación, conciliación y despacho logístico.

## Derechos del titular
El titular puede solicitar acceso, rectificación o supresión de sus datos
escribiendo a privacidad@techretail.com.ar. Respondemos en un plazo máximo
de 10 días corridos.

## Seguridad
Los datos sensibles y credenciales se almacenan cifrados. Las claves de API
se guardan en un gestor de secretos y nunca en texto plano.
""",
        encoding="utf-8",
    )


# ==========================================================================
#  HTML
# ==========================================================================
def gen_html():
    (_dir("comercial") / "faq_soporte.html").write_text(
        """<!DOCTYPE html>
<html lang="es">
<head><meta charset="utf-8"><title>FAQ de Soporte — TechRetail</title></head>
<body>
<h1>Preguntas Frecuentes (FAQ) — TechRetail</h1>

<h2>¿Cómo empiezo a vender?</h2>
<p>Registrás tu comercio, cargás tu catálogo (manual o vía conector con Tango
Gestión / Bejerman) y activás una pasarela de pago. En menos de 24 horas podés
publicar tu tienda.</p>

<h2>¿Qué medios de pago puedo ofrecer?</h2>
<p>Aceptamos MercadoPago y Stripe. Las acreditaciones se concilian
automáticamente contra cada orden de venta.</p>

<h2>¿La facturación es automática?</h2>
<p>Sí. Al aprobarse el pago, el sistema emite la factura electrónica (A, B o C
según tu condición fiscal) contra ARCA y se la envía al comprador. No entregamos
la confirmación de compra sin haber obtenido antes el CAE.</p>

<h2>¿Cuánto tarda el soporte en responder?</h2>
<p>El soporte por chat responde en horario laboral (lunes a viernes de 9 a 18 h)
en un plazo promedio de 2 horas. Los planes Growth y Pro tienen prioridad.</p>

<h2>¿Puedo integrar mi sistema de gestión?</h2>
<p>Sí, ofrecemos conectores nativos con Tango Gestión y Bejerman para
sincronizar catálogo y stock de forma automática.</p>
</body>
</html>
""",
        encoding="utf-8",
    )

    (_dir("rrhh") / "beneficios.html").write_text(
        """<!DOCTYPE html>
<html lang="es"><head><meta charset="utf-8"><title>Beneficios — TechRetail</title></head>
<body>
<h1>Beneficios para Colaboradores — TechRetail</h1>
<ul>
  <li><strong>Vacaciones:</strong> según la Ley de Contrato de Trabajo, más 3 días
  hábiles adicionales por año a partir del segundo año en la empresa.</li>
  <li><strong>Cobertura médica:</strong> prepaga de primer nivel para el colaborador
  y su grupo familiar directo.</li>
  <li><strong>Capacitación:</strong> presupuesto anual de USD 500 para cursos y
  certificaciones.</li>
  <li><strong>Home office:</strong> esquema híbrido y bono único de instalación de
  oficina en casa.</li>
  <li><strong>Día de cumpleaños:</strong> libre.</li>
</ul>
</body>
</html>
""",
        encoding="utf-8",
    )


# ==========================================================================
#  CSV
# ==========================================================================
def gen_csv():
    ruta = _dir("financiero") / "tabla_comisiones_pasarelas.csv"
    filas = [
        ["pasarela", "medio_pago", "comision_porcentaje", "costo_fijo_ars", "plazo_acreditacion_dias"],
        ["MercadoPago", "Tarjeta de crédito", "4.99", "0", "18"],
        ["MercadoPago", "Tarjeta de débito", "3.49", "0", "10"],
        ["MercadoPago", "Dinero en cuenta", "2.99", "0", "0"],
        ["Stripe", "Tarjeta internacional", "3.50", "30", "7"],
        ["Stripe", "Tarjeta local", "2.90", "20", "5"],
    ]
    with ruta.open("w", encoding="utf-8", newline="") as f:
        csv.writer(f).writerows(filas)


# ==========================================================================
#  JSON
# ==========================================================================
def gen_json():
    catalogo_api = {
        "servicio": "TechRetail API",
        "version": "v1",
        "base_url": "https://api.techretail.com.ar/v1",
        "autenticacion": "API Key en el header Authorization: Bearer <token>",
        "endpoints": [
            {
                "ruta": "/productos",
                "metodo": "GET",
                "descripcion": "Lista los productos del catálogo del comercio.",
            },
            {
                "ruta": "/ordenes",
                "metodo": "POST",
                "descripcion": "Crea una orden de compra y dispara el cobro.",
            },
            {
                "ruta": "/facturas/{orden_id}",
                "metodo": "GET",
                "descripcion": "Devuelve la factura electrónica (con CAE) de una orden.",
            },
            {
                "ruta": "/conciliacion",
                "metodo": "GET",
                "descripcion": "Reporte de conciliación entre cobros y órdenes.",
            },
        ],
        "limites": {"rate_limit_por_minuto": 120, "tamano_max_payload_kb": 512},
    }
    (_dir("datos_sistemas") / "catalogo_api.json").write_text(
        json.dumps(catalogo_api, ensure_ascii=False, indent=2), encoding="utf-8"
    )


# ==========================================================================
#  DOCX (Word)
# ==========================================================================
def gen_docx():
    from docx import Document

    # Política de vacaciones (RRHH)
    doc = Document()
    doc.add_heading("Política de Vacaciones y Licencias — TechRetail", level=1)
    doc.add_paragraph(
        "Esta política regula el goce de vacaciones y licencias de todos los "
        "colaboradores de TechRetail Solutions S.R.L."
    )
    doc.add_heading("Vacaciones anuales", level=2)
    doc.add_paragraph(
        "Los días de vacaciones se otorgan según la antigüedad, conforme a la Ley "
        "de Contrato de Trabajo (LCT):"
    )
    tabla = doc.add_table(rows=1, cols=2)
    tabla.style = "Light Grid Accent 1"
    hdr = tabla.rows[0].cells
    hdr[0].text = "Antigüedad"
    hdr[1].text = "Días corridos"
    for antig, dias in [
        ("Menos de 5 años", "14 días"),
        ("De 5 a 10 años", "21 días"),
        ("De 10 a 20 años", "28 días"),
        ("Más de 20 años", "35 días"),
    ]:
        fila = tabla.add_row().cells
        fila[0].text = antig
        fila[1].text = dias
    doc.add_paragraph(
        "Adicionalmente, la empresa otorga 3 días hábiles extra por año a partir "
        "del segundo año de antigüedad."
    )
    doc.add_heading("Licencias especiales", level=2)
    doc.add_paragraph("• Licencia por matrimonio: 10 días corridos.")
    doc.add_paragraph("• Licencia por paternidad: 15 días corridos.")
    doc.add_paragraph("• Licencia por examen: hasta 4 días por año.")
    doc.add_heading("Solicitud", level=2)
    doc.add_paragraph(
        "Las vacaciones se solicitan con al menos 30 días de anticipación a través "
        "del sistema de RH, y deben ser aprobadas por el líder de squad."
    )
    doc.save(str(_dir("rrhh") / "politica_vacaciones_y_licencias.docx"))

    # Procedimiento de conciliación (Operacional)
    doc2 = Document()
    doc2.add_heading("Procedimiento de Conciliación Financiera — TechRetail", level=1)
    doc2.add_paragraph(
        "Describe cómo el Motor de Conciliación Automatizado cruza los cobros de "
        "las pasarelas contra las órdenes de venta."
    )
    doc2.add_heading("Pasos", level=2)
    doc2.add_paragraph("1. La pasarela envía un webhook por cada cobro acreditado.")
    doc2.add_paragraph(
        "2. El motor cruza el monto neto depositado contra la orden, aplicando la "
        "matriz de comisiones configurada por pasarela y medio de pago."
    )
    doc2.add_paragraph(
        "3. Toda orden con una diferencia superior al 0,1% se aísla en el panel de "
        "discrepancias para su resolución inmediata."
    )
    doc2.add_paragraph(
        "El objetivo es eliminar la discrepancia estructural del 2% que existía con "
        "la conciliación manual quincenal en Excel."
    )
    doc2.save(str(_dir("operacional") / "procedimiento_conciliacion.docx"))


# ==========================================================================
#  XLSX (Excel)
# ==========================================================================
def gen_xlsx():
    from openpyxl import Workbook

    # Planes y precios (Comercial)
    wb = Workbook()
    ws = wb.active
    ws.title = "Planes"
    ws.append(["plan", "precio_mensual_ars", "limite_productos", "comercios", "soporte"])
    for fila in [
        ["Starter", 15000, 100, 1, "Email"],
        ["Growth", 39000, 1000, 3, "Chat prioritario"],
        ["Pro", 89000, "Ilimitado", 10, "Chat prioritario + gerente de cuenta"],
    ]:
        ws.append(fila)
    wb.save(str(_dir("comercial") / "planes_y_precios.xlsx"))

    # Modelo de ingresos MRR (Financiero)
    wb2 = Workbook()
    ws2 = wb2.active
    ws2.title = "MRR"
    ws2.append(["mes", "clientes_starter", "clientes_growth", "clientes_pro", "mrr_ars"])
    datos = [
        ["Enero", 120, 45, 8, 3247000],
        ["Febrero", 145, 52, 10, 3903000],
        ["Marzo", 170, 60, 13, 4667000],
        ["Abril", 198, 71, 16, 5573000],
    ]
    for fila in datos:
        ws2.append(fila)
    wb2.save(str(_dir("financiero") / "modelo_ingresos_mrr.xlsx"))


# ==========================================================================
#  PPTX (PowerPoint)
# ==========================================================================
def gen_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # Slide de título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TechRetail Solutions S.R.L."
    slide.placeholders[1].text = "Pitch Comercial — Plataforma SaaS de e-commerce para PyMEs"

    def slide_bullets(titulo, bullets, notas=""):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = titulo
        cuerpo = s.placeholders[1].text_frame
        cuerpo.text = bullets[0]
        for b in bullets[1:]:
            p = cuerpo.add_paragraph()
            p.text = b
            p.level = 0
        if notas:
            s.notes_slide.notes_text_frame.text = notas
        return s

    slide_bullets(
        "El problema",
        [
            "Las PyMEs venden en canales fragmentados.",
            "La facturación fiscal (ARCA) es manual y propensa a errores.",
            "La conciliación de pagos consume ~8 horas semanales por tienda.",
        ],
        notas="Dato clave: la conciliación manual generaba una discrepancia del 2% en caja.",
    )
    slide_bullets(
        "Nuestra solución",
        [
            "Catálogo, cobros, facturación y logística en una sola plataforma.",
            "Facturación automática contra ARCA con CAE en milisegundos.",
            "Motor de conciliación automatizado (discrepancia objetivo < 0,1%).",
        ],
        notas="Integraciones nativas con Tango Gestión y Bejerman.",
    )
    slide_bullets(
        "Planes",
        [
            "Starter: ARS 15.000/mes — ideal para emprendedores.",
            "Growth: ARS 39.000/mes — pensado para comercios en expansión.",
            "Pro: ARS 89.000/mes — para múltiples tiendas y alto volumen.",
        ],
    )
    prs.save(str(_dir("comercial") / "pitch_comercial.pptx"))


# ==========================================================================
#  PDF
# ==========================================================================
def gen_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    estilos = getSampleStyleSheet()

    def construir(ruta, titulo, parrafos):
        doc = SimpleDocTemplate(str(ruta), pagesize=A4)
        flujo = [Paragraph(titulo, estilos["Title"]), Spacer(1, 12)]
        for p in parrafos:
            if p.startswith("## "):
                flujo.append(Paragraph(p[3:], estilos["Heading2"]))
            else:
                flujo.append(Paragraph(p, estilos["BodyText"]))
            flujo.append(Spacer(1, 6))
        doc.build(flujo)

    # Términos y condiciones (Legal)
    construir(
        _dir("legal") / "terminos_y_condiciones.pdf",
        "Términos y Condiciones de Uso — TechRetail",
        [
            "## 1. Aceptación",
            "Al registrarse y utilizar la plataforma TechRetail, el comercio acepta "
            "estos Términos y Condiciones en su totalidad.",
            "## 2. Objeto del servicio",
            "TechRetail provee una plataforma SaaS de e-commerce que incluye catálogo, "
            "cobros, facturación electrónica y conciliación financiera.",
            "## 3. Obligaciones del comercio",
            "El comercio es responsable de la veracidad de sus datos fiscales y del "
            "cumplimiento de sus obligaciones tributarias ante ARCA.",
            "## 4. Facturación del servicio",
            "El plan contratado se factura por mes adelantado. Los precios pueden "
            "ajustarse por indexación (coeficiente UVA) con aviso previo de 30 días.",
            "## 5. Baja del servicio",
            "El comercio puede solicitar la baja en cualquier momento; el servicio "
            "permanece activo hasta el fin del período ya facturado.",
        ],
    )

    # Manual operativo (Operacional)
    construir(
        _dir("operacional") / "manual_procesos_operativos.pdf",
        "Manual de Procesos Operativos — TechRetail",
        [
            "## Proceso de Onboarding de Clientes",
            "El alta de un comercio se realiza mediante un wizard. El flujo está "
            "segmentado: simplificado para emprendedores y avanzado para developers "
            "y agencias. Si se detecta una interrupción en el alta de pagos, el "
            "sistema emite una comunicación por chat a las 4 horas.",
            "## Proceso de Ventas y Facturación",
            "El sistema intercepta la aprobación del pago de forma síncrona y, mediante "
            "el conector con el WSFEv1 de ARCA, emite la factura legal (A, B o C). "
            "No se entrega la confirmación de compra al comprador si no se obtuvo "
            "previamente el CAE.",
            "## Proceso de Gestión de Inventario",
            "El catálogo se sincroniza automáticamente con los ERPs Tango Gestión y "
            "Bejerman mediante conectores nativos, evitando quiebres de stock por "
            "desactualización manual.",
        ],
    )


def main():
    DOCS.mkdir(parents=True, exist_ok=True)
    gen_markdown()
    gen_html()
    gen_csv()
    gen_json()
    gen_docx()
    gen_xlsx()
    gen_pptx()
    gen_pdf()

    archivos = sorted(p for p in DOCS.rglob("*") if p.is_file())
    print(f"Generados {len(archivos)} documentos en {DOCS}:")
    for a in archivos:
        print(f"  - {a.relative_to(DOCS)}")


if __name__ == "__main__":
    main()
