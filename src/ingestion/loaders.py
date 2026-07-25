"""
Extracción de texto por formato.

Cada formato exige un enfoque distinto. La función pública ``extract_text``
recibe una ruta y delega en el extractor correspondiente según la extensión.
Los formatos soportados cubren los más comunes en una empresa:
PDF, Word, Excel, PowerPoint, Markdown, CSV, JSON e HTML.
"""
from __future__ import annotations

import csv
import json
from pathlib import Path

# Extensiones soportadas (para poder recorrer el directorio de documentos).
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx", ".md", ".csv", ".json", ".html", ".htm", ".txt",
}


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    partes = []
    for i, page in enumerate(reader.pages, start=1):
        texto = page.extract_text() or ""
        if texto.strip():
            partes.append(f"[Página {i}]\n{texto}")
    return "\n\n".join(partes)


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    lineas = []
    # Preserva títulos/párrafos para no perder la estructura al hacer chunking.
    for p in doc.paragraphs:
        if p.text.strip():
            lineas.append(p.text)
    # También extrae el contenido de las tablas.
    for tabla in doc.tables:
        for fila in tabla.rows:
            celdas = [c.text.strip() for c in fila.cells]
            if any(celdas):
                lineas.append(" | ".join(celdas))
    return "\n".join(lineas)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True)
    partes = []
    for hoja in wb.worksheets:
        filas = list(hoja.iter_rows(values_only=True))
        if not filas:
            continue
        encabezados = [str(c) if c is not None else "" for c in filas[0]]
        partes.append(f"[Hoja: {hoja.title}]")
        # Convierte cada fila en texto legible repitiendo los encabezados,
        # que es la forma recomendada de representar planillas para RAG.
        for fila in filas[1:]:
            pares = [
                f"{h}: {v}"
                for h, v in zip(encabezados, fila)
                if v is not None and str(v).strip()
            ]
            if pares:
                partes.append("; ".join(pares))
    return "\n".join(partes)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    partes = []
    for i, slide in enumerate(prs.slides, start=1):
        textos = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                textos.append(shape.text_frame.text)
        # Notas del orador: suelen tener contexto adicional valioso.
        if slide.has_notes_slide:
            nota = slide.notes_slide.notes_text_frame.text
            if nota.strip():
                textos.append(f"(Notas) {nota}")
        if textos:
            partes.append(f"[Diapositiva {i}]\n" + "\n".join(textos))
    return "\n\n".join(partes)


def _extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    # Elimina ruido técnico que no aporta significado.
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _extract_csv(path: Path) -> str:
    partes = []
    with path.open(encoding="utf-8", errors="ignore", newline="") as f:
        lector = csv.reader(f)
        filas = list(lector)
    if not filas:
        return ""
    encabezados = filas[0]
    for fila in filas[1:]:
        pares = [f"{h}: {v}" for h, v in zip(encabezados, fila) if v.strip()]
        if pares:
            partes.append("; ".join(pares))
    return "\n".join(partes)


def _extract_json(path: Path) -> str:
    datos = json.loads(path.read_text(encoding="utf-8", errors="ignore"))

    def aplanar(obj, prefijo=""):
        lineas = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                lineas.extend(aplanar(v, f"{prefijo}{k}."))
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                lineas.extend(aplanar(v, f"{prefijo}{i}."))
        else:
            lineas.append(f"{prefijo.rstrip('.')}: {obj}")
        return lineas

    return "\n".join(aplanar(datos))


def _extract_plain(path: Path) -> str:
    # Markdown, TXT y similares: texto ya legible.
    return path.read_text(encoding="utf-8", errors="ignore")


_EXTRACTORES = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".csv": _extract_csv,
    ".json": _extract_json,
    ".md": _extract_plain,
    ".txt": _extract_plain,
}


def extract_text(path: Path) -> str:
    """Extrae el texto de un archivo según su extensión.

    Lanza ``ValueError`` si el formato no está soportado.
    """
    ext = path.suffix.lower()
    extractor = _EXTRACTORES.get(ext)
    if extractor is None:
        raise ValueError(f"Formato no soportado: {ext} ({path.name})")
    return extractor(path)
